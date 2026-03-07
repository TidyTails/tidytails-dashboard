#!/usr/bin/env python3
"""
MOTRA Pitch Deck v6 - NO OVERLAPPING SHAPES
Key fix: All text is either:
1. Standalone textbox (no background)
2. Or embedded IN the shape's text_frame (no separate textbox)

NO overlapping shapes at all.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE = RgbColor(0, 102, 255)
DARK = RgbColor(26, 26, 26)
GRAY = RgbColor(102, 102, 102)
WHITE = RgbColor(255, 255, 255)
RED = RgbColor(220, 38, 38)
GREEN = RgbColor(22, 163, 74)
DARK_BLUE = RgbColor(0, 40, 80)
CARD_BG = RgbColor(30, 30, 30)

def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

def txt(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Simple textbox - no background"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align

def box_with_text(slide, left, top, width, height, text, text_size=14, text_color=WHITE, bg_color=DARK_BLUE, border_color=BLUE, bold=False, align=PP_ALIGN.CENTER):
    """Shape with text INSIDE it - no overlapping"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    
    # Put text IN the shape
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(text_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.alignment = align
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    
    return shape

# === SLIDE 0: Title ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 2.8, 12.333, 1, "MOTRA", size=72, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
txt(slide, 0.5, 4.2, 12.333, 0.5, "AUTONOMY, MAINTAINED", size=24, color=GRAY, align=PP_ALIGN.CENTER)
txt(slide, 0.5, 6.5, 12.333, 0.4, "Seed Round — March 2026", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# === SLIDE 1: Company Purpose ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "01 / COMPANY PURPOSE", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "The infrastructure layer for AV fleet care.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.9, 12, 0.8, "MOTRA is the infrastructure layer for autonomous vehicle fleet care.", size=24, color=WHITE)
txt(slide, 0.5, 3.0, 12, 1, "We deploy a gig-powered network of mobile technicians to clean and service robotaxis on-location, on-demand — keeping fleets running 24/7.", size=16, color=GRAY)
box_with_text(slide, 0.5, 4.5, 12, 1, "Think DoorDash for fleet maintenance — invisible, essential, everywhere.", text_size=18, bold=True)

# === SLIDE 2: Problem ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "02 / THE PROBLEM", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "The maintenance bottleneck.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.8, 12, 0.6, "Robotaxis run 15+ hours a day. Their maintenance infrastructure doesn't.", size=20, color=WHITE)

txt(slide, 0.5, 2.6, 5, 0.35, "1. No driver = no eyes", size=14, color=WHITE, bold=True)
txt(slide, 0.5, 2.95, 5, 0.3, "Nobody notices trash, spills, or wear", size=11, color=GRAY)
txt(slide, 0.5, 3.35, 5, 0.35, "2. Depot-dependent", size=14, color=WHITE, bold=True)
txt(slide, 0.5, 3.7, 5, 0.3, "Vehicles must return for basic cleaning", size=11, color=GRAY)
txt(slide, 0.5, 4.1, 5, 0.35, "3. Downtime = lost revenue", size=14, color=WHITE, bold=True)
txt(slide, 0.5, 4.45, 5, 0.3, "Every minute in depot is a ride not taken", size=11, color=GRAY)
txt(slide, 0.5, 4.85, 5, 0.35, "4. Sensors are critical", size=14, color=WHITE, bold=True)
txt(slide, 0.5, 5.2, 5, 0.3, "Dirty LiDAR/cameras = degraded driving", size=11, color=GRAY)

box_with_text(slide, 6.5, 2.6, 6, 2.8, "Current Reality\n\nWaymo uses 1,000+ Transdev operators at centralized depots.\n\nThat doesn't scale to 1 million rides/week.", text_size=14, bg_color=RgbColor(50, 30, 30), border_color=RED)

# === SLIDE 3: Solution ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "03 / THE SOLUTION", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "Mobile fleet care on demand.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.8, 12, 0.5, "MOTRA deploys certified technicians directly to vehicles — wherever they are.", size=16, color=GRAY)

box_with_text(slide, 0.5, 2.6, 3.9, 2.2, "Quick Clean\n$12-18\n5-10 min\nWipe-down, trash, odor", text_size=16, bg_color=CARD_BG, border_color=RgbColor(60,60,60))
box_with_text(slide, 4.6, 2.6, 3.9, 2.2, "Deep Clean\n$45-75\n30-60 min\nFull detail, sensors, exterior", text_size=16, bg_color=CARD_BG, border_color=RgbColor(60,60,60))
box_with_text(slide, 8.7, 2.6, 3.9, 2.2, "Emergency\n$75-150\nOn-demand\nSpills, incidents, urgent", text_size=16, bg_color=CARD_BG, border_color=RgbColor(60,60,60))

box_with_text(slide, 0.5, 5.2, 12, 0.8, "Vehicles stay in zone  |  Variable cost  |  Scales instantly  |  24/7 availability", text_size=14)

# === SLIDE 4: Why Now ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "04 / WHY NOW", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "Escape velocity.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.7, 12, 0.5, "The robotaxi industry just hit an inflection point.", size=18, color=GRAY)

box_with_text(slide, 0.5, 2.4, 2.9, 1.3, "400K\nRides/Week", text_size=24, text_color=BLUE, bold=True)
box_with_text(slide, 3.6, 2.4, 2.9, 1.3, "10\nUS Markets", text_size=24, text_color=BLUE, bold=True)
box_with_text(slide, 6.7, 2.4, 2.9, 1.3, "$16B\nFresh Funding", text_size=24, text_color=BLUE, bold=True)
box_with_text(slide, 9.8, 2.4, 2.9, 1.3, "$126B\nValuation", text_size=24, text_color=BLUE, bold=True)

txt(slide, 0.5, 4.0, 2.5, 0.35, "Metric", size=11, color=BLUE, bold=True)
txt(slide, 3.2, 4.0, 3.5, 0.35, "Now (Feb 2026)", size=11, color=BLUE, bold=True)
txt(slide, 7, 4.0, 5, 0.35, "EOY 2026 Target", size=11, color=BLUE, bold=True)

rows = [("Weekly Rides", "400,000", "1,000,000"), ("Markets", "10 cities", "20+ cities (London, Tokyo)"), ("Fleet Size", "2,500 vehicles", "4,500+ vehicles"), ("Price War", "Waymo $8.17 vs Uber $17.25", "Volume exploding")]
y = 4.4
for metric, now, target in rows:
    txt(slide, 0.5, y, 2.5, 0.35, metric, size=12, color=WHITE)
    txt(slide, 3.2, y, 3.5, 0.35, now, size=12, color=WHITE, bold=True)
    txt(slide, 7, y, 5, 0.35, target, size=12, color=WHITE)
    y += 0.4
txt(slide, 0.5, 6.2, 12, 0.3, "Sources: TechXplore Feb 2026, Fox News Feb 2026, Bloomberg", size=9, color=GRAY)

# === SLIDE 5: Market Size ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "05 / MARKET SIZE", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "A $147B opportunity.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.7, 12, 0.5, "Robotaxi market: 99% CAGR through 2033", size=18, color=GRAY)

box_with_text(slide, 0.5, 2.4, 3.8, 1.3, "$2.7B\n2025 TAM", text_size=28, text_color=BLUE, bold=True)
box_with_text(slide, 4.6, 2.4, 3.8, 1.3, "$147B\n2033 TAM", text_size=28, text_color=BLUE, bold=True)
box_with_text(slide, 8.7, 2.4, 3.8, 1.3, "99%\nCAGR", text_size=28, text_color=BLUE, bold=True)

txt(slide, 0.5, 4.1, 12, 0.4, "MOTRA's Serviceable Market", size=16, color=BLUE, bold=True)
txt(slide, 0.5, 4.6, 1.2, 0.35, "2028", size=14, color=WHITE, bold=True)
txt(slide, 2, 4.6, 1.8, 0.35, "$547M", size=14, color=BLUE, bold=True)
txt(slide, 4.2, 4.6, 8, 0.35, "100K robotaxis × ~$105/week", size=12, color=GRAY)
txt(slide, 0.5, 5.0, 1.2, 0.35, "2030", size=14, color=WHITE, bold=True)
txt(slide, 2, 5.0, 1.8, 0.35, "$2.2B", size=14, color=BLUE, bold=True)
txt(slide, 4.2, 5.0, 8, 0.35, "200K robotaxis", size=12, color=GRAY)
txt(slide, 0.5, 5.4, 1.2, 0.35, "2032", size=14, color=WHITE, bold=True)
txt(slide, 2, 5.4, 1.8, 0.35, "$5.5B", size=14, color=BLUE, bold=True)
txt(slide, 4.2, 5.4, 8, 0.35, "500K robotaxis + commercial EVs", size=12, color=GRAY)
txt(slide, 0.5, 6.2, 12, 0.3, "Source: Grand View Research, MOTRA analysis", size=9, color=GRAY)

# === SLIDE 6: Competition ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "06 / COMPETITION", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "Our competitive edge.", size=32, color=WHITE, bold=True)

# PURE TEXT TABLE - no shapes
txt(slide, 0.5, 1.9, 3, 0.35, "Player", size=11, color=BLUE, bold=True)
txt(slide, 4, 1.9, 4, 0.35, "Approach", size=11, color=BLUE, bold=True)
txt(slide, 8.5, 1.9, 4, 0.35, "Limitation", size=11, color=BLUE, bold=True)

rows = [("Uber AV Services", "Depot-based cleaning + charging", "Competitor to Waymo"), ("Transdev", "1,000+ operators at depots", "Depot-bound, not mobile"), ("Tesla Robots", "Automated cleaning (patents)", "Years away, Tesla-only"), ("In-House Ops", "AV companies DIY", "High fixed cost")]
y = 2.4
for player, approach, limitation in rows:
    txt(slide, 0.5, y, 3.3, 0.4, player, size=13, color=WHITE, bold=True)
    txt(slide, 4, y, 4, 0.4, approach, size=11, color=GRAY)
    txt(slide, 8.5, y, 4, 0.4, limitation, size=11, color=RED)
    y += 0.5

box_with_text(slide, 0.5, 4.8, 12, 1, "MOTRA Advantage: Neutral (Switzerland) · Mobile-first · 100% focused on AV fleet care", text_size=16, text_color=GREEN, bg_color=RgbColor(20, 50, 30), border_color=GREEN, bold=True)

# === SLIDE 7: Why Not In-House ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "07 / WHY NOT IN-HOUSE", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "The outsourcing thesis.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.7, 12, 0.5, "Why won't Waymo just do this themselves?", size=18, color=GRAY)

box_with_text(slide, 0.5, 2.4, 3.9, 2.4, "1. Already Outsource\n\nWaymo uses Transdev, Moove, Avis, TechForce, Amerit, Terawatt today.\n\nBuild vs. buy = they buy.", text_size=11, bg_color=CARD_BG, border_color=BLUE, align=PP_ALIGN.LEFT)
box_with_text(slide, 4.6, 2.4, 3.9, 2.4, "2. Not Core Competency\n\nWaymo's edge is autonomous driving AI.\n\nEvery $ on cleaning = $ not on AI.", text_size=11, bg_color=CARD_BG, border_color=BLUE, align=PP_ALIGN.LEFT)
box_with_text(slide, 8.7, 2.4, 3.9, 2.4, "3. Math Doesn't Work\n\n20+ cities = thousands of employees, HR, facilities.\n\nGig model = variable cost.", text_size=11, bg_color=CARD_BG, border_color=BLUE, align=PP_ALIGN.LEFT)

box_with_text(slide, 0.5, 5.2, 12, 0.9, "Uber just validated this — if big players could do it in-house, they wouldn't build outsourced solutions.", text_size=14)

# === SLIDE 8: Moat ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "08 / DEFENSIBILITY", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "Our moat.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.7, 12, 0.5, "Three layers of defensibility.", size=18, color=GRAY)

box_with_text(slide, 0.5, 2.4, 3.9, 2.8, "Year 1-2: Network Effects\n\n• First AV gig network per city\n• Hyperlocal network effects\n• More techs → faster response → more volume", text_size=10, bg_color=CARD_BG, border_color=BLUE, align=PP_ALIGN.LEFT)
box_with_text(slide, 4.6, 2.4, 3.9, 2.8, "Year 2-3: Data Advantage\n\n• Every service = operational data\n• Predictive maintenance insights\n• DoorDash: 30% of eng in data roles", text_size=10, bg_color=CARD_BG, border_color=BLUE, align=PP_ALIGN.LEFT)
box_with_text(slide, 8.7, 2.4, 3.9, 2.8, "Year 3+: Integration Lock-In\n\n• API embedded in fleet systems\n• Part of operator workflows\n• 5% retention = 25-95% profit boost", text_size=10, bg_color=CARD_BG, border_color=BLUE, align=PP_ALIGN.LEFT)

txt(slide, 0.5, 5.5, 12, 0.4, "The Uber analogy: Uber's moat isn't the app — it's the network of drivers in every city.", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# === SLIDE 9: Team ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "09 / TEAM", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "Leadership.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.7, 12, 0.5, "15 years at Boeing. 2 patents. Built to scale.", size=18, color=GRAY)

box_with_text(slide, 0.5, 2.4, 6, 3.2, "Adonis Williams — Founder & CEO\n\nRole: Deputy Functional Chief Engineer, Boeing\nTenure: 15 years at Boeing\nPatents: 2 (cybersecurity, blockchain comms)\nEducation: MS Systems Engineering, Missouri S&T", text_size=12, bg_color=CARD_BG, border_color=BLUE, align=PP_ALIGN.LEFT)
box_with_text(slide, 7, 2.4, 5.5, 3.2, "Why Adonis Wins This\n\nMOTRA is a systems problem: coordinating a distributed workforce to service safety-critical AVs.\n\nThat's exactly what he's done at Boeing for 15 years.\n\n✓ Complex operations at scale\n✓ Distributed teams\n✓ Safety-critical systems", text_size=11, bg_color=RgbColor(20, 25, 35), border_color=RgbColor(20, 25, 35), align=PP_ALIGN.LEFT)

# === SLIDE 10: Risk Mitigation ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "10 / RISK MITIGATION", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "De-risked.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.7, 12, 0.5, "What if the AV market slows down?", size=18, color=GRAY)

box_with_text(slide, 0.5, 2.4, 5.8, 1.4, "The Risk\n\nRobotaxi adoption could hit regulatory delays, technical hurdles, or economic headwinds.", text_size=11, bg_color=RgbColor(50, 30, 30), border_color=RED, align=PP_ALIGN.LEFT)
box_with_text(slide, 6.8, 2.4, 5.8, 1.4, "Adjacent Markets\n\nEven if AV slows, 250K+ commercial EVs need service.", text_size=11, bg_color=RgbColor(20, 50, 30), border_color=GREEN, align=PP_ALIGN.LEFT)

txt(slide, 0.5, 4.1, 3, 0.35, "Fleet", size=11, color=BLUE, bold=True)
txt(slide, 4, 4.1, 3, 0.35, "Current", size=11, color=BLUE, bold=True)
txt(slide, 8, 4.1, 4, 0.35, "Target", size=11, color=BLUE, bold=True)

fleets = [("Amazon (Rivian)", "30,000 EVs", "100,000 by 2030"), ("FedEx", "200,000+ vehicles", "All-electric by 2040"), ("UPS", "18,300+ alt-fuel", "10,000 more ordered")]
y = 4.5
for fleet, current, target in fleets:
    txt(slide, 0.5, y, 3.3, 0.35, fleet, size=12, color=WHITE, bold=True)
    txt(slide, 4, y, 3.5, 0.35, current, size=11, color=WHITE)
    txt(slide, 8, y, 4, 0.35, target, size=11, color=WHITE)
    y += 0.4

txt(slide, 0.5, 5.9, 12, 0.4, "MOTRA has multiple paths to scale. Robotaxis are the beachhead — commercial EVs are the expansion.", size=14, color=GREEN, align=PP_ALIGN.CENTER)

# === SLIDE 11: GTM ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "11 / GO-TO-MARKET", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "First customers.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.7, 12, 0.5, "The DoorDash playbook: do things that don't scale.", size=18, color=GRAY)

box_with_text(slide, 0.5, 2.4, 3.9, 2.8, "Phase 1: Manual Pilot\nMonth 1-3\n\n• ONE operator, ONE city\n• Free 30-day, 50 vehicles\n• Text dispatch, Google Sheets\n• We ARE the operators", text_size=10, bg_color=CARD_BG, border_color=BLUE, align=PP_ALIGN.LEFT)
box_with_text(slide, 4.6, 2.4, 3.9, 2.8, "Phase 2: Paid Contract\nMonth 4-6\n\n• Convert to paid\n• Price 20% below Transdev\n• 6-month commitment\n• Build reference customer", text_size=10, bg_color=CARD_BG, border_color=BLUE, align=PP_ALIGN.LEFT)
box_with_text(slide, 8.7, 2.4, 3.9, 2.8, "Phase 3: Scale\nMonth 7-12\n\n• Add second city\n• Launch MVP app\n• First customer as reference\n• Target $100K MRR", text_size=10, bg_color=CARD_BG, border_color=BLUE, align=PP_ALIGN.LEFT)

txt(slide, 0.5, 5.5, 12, 0.4, "Warm intros = 3-5x better close rate. We're asking them to try a service, not trust an app.", size=14, color=WHITE, align=PP_ALIGN.CENTER)

# === SLIDE 12: Business Model ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "12 / BUSINESS MODEL", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "Unit economics.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.7, 12, 0.5, "Platform take-rate on every service.", size=18, color=GRAY)

txt(slide, 0.5, 2.4, 3.5, 0.4, "Avg service price", size=13, color=WHITE)
txt(slide, 4.5, 2.4, 2.5, 0.4, "$15", size=13, color=BLUE, bold=True)
txt(slide, 0.5, 2.85, 3.5, 0.4, "Tech payout", size=13, color=WHITE)
txt(slide, 4.5, 2.85, 2.5, 0.4, "$10-11 (65-70%)", size=13, color=BLUE, bold=True)
txt(slide, 0.5, 3.3, 3.5, 0.4, "Platform margin", size=13, color=WHITE)
txt(slide, 4.5, 3.3, 2.5, 0.4, "$4-5 (27-33%)", size=13, color=BLUE, bold=True)
txt(slide, 0.5, 3.75, 3.5, 0.4, "Services/tech/day", size=13, color=WHITE)
txt(slide, 4.5, 3.75, 2.5, 0.4, "15-20", size=13, color=BLUE, bold=True)
txt(slide, 0.5, 4.2, 3.5, 0.4, "Tech daily earnings", size=13, color=WHITE)
txt(slide, 4.5, 4.2, 2.5, 0.4, "$150-220", size=13, color=BLUE, bold=True)

box_with_text(slide, 8, 2.4, 4, 1.8, "27-33%\nPlatform Margin\n\nHigh operating leverage — costs don't scale linearly", text_size=12, text_color=BLUE, bold=True)

box_with_text(slide, 0.5, 5.0, 12, 0.9, "At Scale: 10,000 services/day = $150K revenue, $40-50K margin/day", text_size=16, bold=True)

# === SLIDE 13: The Ask ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "13 / THE ASK", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "Investment.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.7, 12, 0.5, "$1.5M Seed to prove the model.", size=18, color=GRAY)

box_with_text(slide, 0.5, 2.4, 3.5, 1.3, "$1.5M\nSeed Round", text_size=28, text_color=BLUE, bold=True)

txt(slide, 4.5, 2.4, 3, 0.4, "Use of Funds", size=14, color=BLUE, bold=True)
txt(slide, 4.5, 2.9, 2.2, 0.35, "Product (40%)", size=11, color=WHITE, bold=True)
txt(slide, 6.8, 2.9, 3, 0.35, "MVP app, API integrations", size=10, color=GRAY)
txt(slide, 4.5, 3.3, 2.2, 0.35, "Operations (30%)", size=11, color=WHITE, bold=True)
txt(slide, 6.8, 3.3, 3, 0.35, "Techs, training, equipment", size=10, color=GRAY)
txt(slide, 4.5, 3.7, 2.2, 0.35, "Sales (20%)", size=11, color=WHITE, bold=True)
txt(slide, 6.8, 3.7, 3, 0.35, "Enterprise BD", size=10, color=GRAY)
txt(slide, 4.5, 4.1, 2.2, 0.35, "G&A (10%)", size=11, color=WHITE, bold=True)
txt(slide, 6.8, 4.1, 3, 0.35, "Legal, insurance", size=10, color=GRAY)

txt(slide, 0.5, 4.5, 12, 0.4, "18-Month Milestones", size=14, color=BLUE, bold=True)
txt(slide, 0.5, 5.0, 3.5, 0.35, "First Paid Contract", size=12, color=WHITE)
txt(slide, 4.5, 5.0, 2, 0.35, "Month 6", size=12, color=GRAY)
txt(slide, 7, 5.0, 4, 0.35, "$10K+ MRR", size=12, color=GREEN, bold=True)
txt(slide, 0.5, 5.4, 3.5, 0.35, "Product-Market Fit", size=12, color=WHITE)
txt(slide, 4.5, 5.4, 2, 0.35, "Month 12", size=12, color=GRAY)
txt(slide, 7, 5.4, 4, 0.35, "$100K MRR, 3 cities", size=12, color=GREEN, bold=True)
txt(slide, 0.5, 5.8, 3.5, 0.35, "Series A Ready", size=12, color=WHITE)
txt(slide, 4.5, 5.8, 2, 0.35, "Month 18", size=12, color=GRAY)
txt(slide, 7, 5.8, 4, 0.35, "$300K MRR, 1,000+ vehicles", size=12, color=GREEN, bold=True)

# === SLIDE 14: Exit ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 0.4, 4, 0.3, "14 / EXIT STRATEGY", size=10, color=BLUE, bold=True)
txt(slide, 0.5, 0.9, 12, 0.8, "Returns.", size=32, color=WHITE, bold=True)
txt(slide, 0.5, 1.7, 12, 0.5, "Three paths to liquidity.", size=18, color=GRAY)

box_with_text(slide, 0.5, 2.4, 3.9, 2.8, "Acquisition (Most Likely)\n3-5 years, $100-500M\n\n• Waymo (uses 6+ contractors)\n• Uber (just launched AV Services)\n• Amazon (owns Zoox)\n\nComp: Zoox = $1.2B", text_size=10, bg_color=CARD_BG, border_color=GREEN, align=PP_ALIGN.LEFT)
box_with_text(slide, 4.6, 2.4, 3.9, 2.8, "Fleet Services Roll-Up\n4-6 years\n\n• Enterprise Mobility\n• Cox Automotive\n• Fleetio ($1.5B valuation)", text_size=10, bg_color=CARD_BG, border_color=RgbColor(60,60,60), align=PP_ALIGN.LEFT)
box_with_text(slide, 8.7, 2.4, 3.9, 2.8, "IPO Path\n7-10 years\n\n• Expand beyond AV to ALL\n  commercial fleet services\n• 250K+ vehicles addressable", text_size=10, bg_color=CARD_BG, border_color=RgbColor(60,60,60), align=PP_ALIGN.LEFT)

txt(slide, 0.5, 5.5, 12, 0.4, "Waymo is valued at $126B. A $200M acquisition is a rounding error.", size=14, color=WHITE, align=PP_ALIGN.CENTER)

# === SLIDE 15: Close ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
txt(slide, 0.5, 2.2, 12.333, 0.8, "The future of mobility needs infrastructure.", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(slide, 0.5, 3.2, 12.333, 0.5, "AUTONOMY, MAINTAINED.", size=20, color=GRAY, align=PP_ALIGN.CENTER)
txt(slide, 0.5, 4.5, 12.333, 0.4, "Adonis Williams — Founder & CEO", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(slide, 0.5, 5.0, 12.333, 0.35, "Deputy Functional Chief Engineer, Boeing (15 years)", size=13, color=GRAY, align=PP_ALIGN.CENTER)
txt(slide, 0.5, 5.4, 12.333, 0.35, "MS Systems Engineering | 2 Patents", size=13, color=GRAY, align=PP_ALIGN.CENTER)
txt(slide, 0.5, 6.0, 12.333, 0.35, "adonis@motra.io", size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# Save
prs.save('/Users/halpininc/.openclaw/workspace/motra-dashboard/pitch/MOTRA-Pitch-Deck-v6.pptx')
print("Saved v6")
